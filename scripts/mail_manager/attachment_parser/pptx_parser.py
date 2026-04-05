"""
PowerPoint parser implementation using python-pptx.

Extracts text and metadata from PowerPoint presentations (.pptx).
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation


class PPTXParser:
    """Parser for PowerPoint documents using python-pptx."""

    def can_parse(self, file_path: Path) -> bool:
        """
        Check if this parser can handle PowerPoint files.

        Args:
            file_path: Path to the file to check.

        Returns:
            True if file has .pptx extension.
        """
        return file_path.suffix.lower() == ".pptx"

    def extract_text(self, file_path: Path) -> str:
        """
        Extract text content from PowerPoint file.

        Iterates through all slides and all shapes, extracting text
        from shapes that have text content.

        Args:
            file_path: Path to the PowerPoint file.

        Returns:
            Extracted text from all slides, joined by newlines.
        """
        prs = Presentation(str(file_path))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)

    def extract_metadata(self, file_path: Path) -> dict[str, Any]:
        """
        Extract metadata from PowerPoint file.

        Args:
            file_path: Path to the PowerPoint file.

        Returns:
            Dictionary with slide_count and type.
        """
        prs = Presentation(str(file_path))
        return {
            "slide_count": len(prs.slides),
            "type": "powerpoint",
        }
